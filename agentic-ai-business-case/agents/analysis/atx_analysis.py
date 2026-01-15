import os
import pandas as pd
from strands import Agent, tool
from strands.models import BedrockModel
from pypdf import PdfReader
from pptx import Presentation

from agents.config.config import input_folder_dir_path, model_id_claude3_7, model_temperature


# Create a BedrockModel
bedrock_model = BedrockModel(
    model_id=model_id_claude3_7,
    temperature=model_temperature
)

def read_file_from_input_dir(filename):
    """Read file from the input directory"""
    from agents.utils.project_context import get_input_file_path
    return get_input_file_path(filename)

@tool(name="read_excel_file", description="Read Excel file from the input folder and return its content as a dataframe")
def read_excel_file(filename: str):
    """Read Excel file and return its content"""
    full_path = read_file_from_input_dir(filename)
    df = pd.read_excel(full_path, sheet_name=None)  # Read all sheets
    return df

@tool(name="read_pdf_file", description="Read PDF file from the input folder and extract text content")
def read_pdf_file(filename: str):
    """Read PDF file and extract text content"""
    full_path = read_file_from_input_dir(filename)
    reader = PdfReader(full_path)
    text_content = []
    for page_num, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        text_content.append(f"--- Page {page_num} ---\n{text}")
    return "\n\n".join(text_content)

@tool(name="read_pptx_file", description="Read PowerPoint file from the input folder and extract text content from slides")
def read_pptx_file(filename: str):
    """Read PowerPoint file and extract text content from slides"""
    full_path = read_file_from_input_dir(filename)
    prs = Presentation(full_path)
    slide_content = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        slide_text.append(f"--- Slide {slide_num} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slide_content.append("\n".join(slide_text))
    return "\n\n".join(slide_content)

# System message for the ATX analysis agent
system_message = """
You are an expert AWS migration specialist with deep expertise in VMware to AWS cloud migrations. 
Your role is to analyze AWS Transform for VMware (ATX) assessment outputs to extract critical insights 
for building comprehensive business cases for VMware workload migration to AWS.

**About ATX**: AWS Transform for VMware analyzes VMware environments and generates 
detailed reports to help plan and execute migrations from VMware to AWS.

You have access to three ATX output files:
1. **analysis.xlsx** - Excel file containing quantitative VMware environment data, resource utilization, and cost analysis
2. **business_case.pptx** - PowerPoint presentation with executive summary and business case justification
3. **report.pdf** - Comprehensive technical assessment report with detailed VMware infrastructure analysis

When analyzing these ATX documents, focus on:

## (1) VMware Environment Overview
- Extract VMware infrastructure inventory (vCPUs, memory, storage, VMs count)
- Identify VMware versions, clusters, and datacenter configuration
- Document current VMware licensing and support costs
- Assess overall environment complexity and scale

## (2) Workload Analysis & Categorization
- Identify workload types and their characteristics
- Categorize VMs by migration readiness (easy, moderate, complex)
- Document application dependencies and groupings
- Assess workload performance requirements and patterns

## (3) AWS Target Architecture & Mapping
- Extract recommended AWS services for VMware workloads (EC2, Amazon Elastic VMware Service (EVS), etc.)
- Document instance type recommendations and rightsizing opportunities
- Identify modernization opportunities (containers, serverless, managed services)
- Review network architecture and connectivity requirements

## (4) Cost Analysis & TCO Comparison
- Extract current VMware infrastructure costs (hardware, licensing, maintenance, facilities)
- Document projected AWS costs (compute, storage, data transfer, support)
- Analyze cost optimization opportunities and savings
- Review TCO comparison over 3-5 year period
- Identify cost drivers and optimization levers

## (5) Migration Strategy & Approach
- Extract recommended migration patterns (rehost, replatform, refactor)
- Document migration waves and prioritization
- Identify pilot candidates and quick wins
- Review migration timeline and phases

## (6) Risk Assessment & Readiness
- Identify technical risks and blockers
- Document application compatibility issues
- Assess organizational readiness and skill gaps
- Review compliance and security considerations

## (7) Business Case & ROI
- Extract financial benefits and cost savings
- Document operational benefits (agility, scalability, reliability)
- Review ROI projections and payback period
- Identify strategic business value and innovation opportunities

## (8) Key Recommendations & Next Steps
- Summarize strategic recommendations for VMware to AWS migration
- Prioritize immediate actions and quick wins
- Document long-term transformation roadmap
- Highlight critical success factors

**IMPORTANT**: Base your analysis strictly on the ATX assessment data found in the provided documents. 
Do not make assumptions or add information not present in the ATX outputs. Focus on VMware-specific 
metrics and AWS migration considerations.

Format your response in markdown with clear headings, bullet points, and tables where appropriate.
"""

# Create the agent with all three tools
agent = Agent(
    model=bedrock_model,
    system_prompt=system_message,
    tools=[read_excel_file, read_pdf_file, read_pptx_file]
)

# Example usage (commented out)
# question = """
# Analyze the AWS Transform for VMware (ATX) assessment outputs available in the input folder:
# - analysis.xlsx (VMware environment data and cost analysis)
# - business_case.pptx (Executive business case presentation)
# - report.pdf (Detailed technical assessment report)
# 
# Provide a comprehensive analysis of the VMware environment and AWS migration recommendations 
# covering all key aspects for building a business case for VMware to AWS migration.
# """
# 
# result = agent(question)
# print(result.message)


if __name__ == "__main__":
    # Test the ATX analysis tools individually
    print("Testing AWS Transform for VMware (ATX) Analysis Tools...")
    
    # Test Excel reading (VMware environment data)
    try:
        excel_data = read_excel_file("analysis.xlsx")
        print(f"\n✓ ATX analysis.xlsx loaded successfully with {len(excel_data)} sheet(s)")
    except Exception as e:
        print(f"\n✗ Error reading ATX analysis.xlsx: {e}")
    
    # Test PDF reading (Technical assessment report)
    try:
        pdf_content = read_pdf_file("report.pdf")
        print(f"✓ ATX report.pdf loaded successfully ({len(pdf_content)} characters)")
    except Exception as e:
        print(f"✗ Error reading ATX report.pdf: {e}")
    
    # Test PowerPoint reading (Business case presentation)
    try:
        pptx_content = read_pptx_file("business_case.pptx")
        print(f"✓ ATX business_case.pptx loaded successfully ({len(pptx_content)} characters)")
    except Exception as e:
        print(f"✗ Error reading ATX business_case.pptx: {e}")
