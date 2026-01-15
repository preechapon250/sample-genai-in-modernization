"""
ATX PowerPoint Extractor
Extracts key information from ATX Business Case PowerPoint presentation
- Assessment Scope slide
- Executive Summary slide
- Financial Overview slide
"""

from pptx import Presentation
import re
from typing import Dict, Optional


def extract_text_from_slide(slide) -> str:
    """Extract all text from a slide."""
    text_parts = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text_parts.append(shape.text)
    return "\n".join(text_parts)


def find_slide_by_title(prs: Presentation, title_keywords: list) -> Optional[int]:
    """Find slide index by matching title keywords."""
    for idx, slide in enumerate(prs.slides):
        slide_text = extract_text_from_slide(slide).lower()
        if any(keyword.lower() in slide_text for keyword in title_keywords):
            return idx
    return None


def extract_assessment_scope(prs: Presentation) -> Dict:
    """
    Extract Assessment Scope information from ATX PowerPoint.
    
    Looks for slides with titles like:
    - "Assessment Scope"
    - "Scope"
    - "Project Scope"
    """
    scope_keywords = ["assessment scope", "project scope", "scope of assessment"]
    slide_idx = find_slide_by_title(prs, scope_keywords)
    
    if slide_idx is None:
        return {"found": False, "content": ""}
    
    slide = prs.slides[slide_idx]
    content = extract_text_from_slide(slide)
    
    # Parse key metrics
    result = {
        "found": True,
        "content": content,
        "vm_count": None,
        "vcpu_count": None,
        "ram_gb": None,
        "storage_tb": None,
        "storage_tib": None,
        "windows_vms": None,
        "linux_vms": None,
        "database_count": None
    }
    
    # Extract numbers using regex - ATX specific patterns
    # Total servers: 368
    vm_match = re.search(r'total\s+servers?:\s*(\d+)', content, re.IGNORECASE)
    if vm_match:
        result["vm_count"] = int(vm_match.group(1))
    
    # Windows Servers in scope: 16
    windows_match = re.search(r'windows\s+servers?\s+in\s+scope:\s*(\d+)', content, re.IGNORECASE)
    if windows_match:
        result["windows_vms"] = int(windows_match.group(1))
    
    # Linux Servers in scope: 352
    linux_match = re.search(r'linux\s+servers?\s+in\s+scope:\s*(\d+)', content, re.IGNORECASE)
    if linux_match:
        result["linux_vms"] = int(linux_match.group(1))
    
    # Provisioned Storage (TiB): 12
    storage_tib_match = re.search(r'provisioned\s+storage\s+\(tib\):\s*(\d+)', content, re.IGNORECASE)
    if storage_tib_match:
        result["storage_tib"] = float(storage_tib_match.group(1))
        result["storage_tb"] = float(storage_tib_match.group(1)) * 1.09951  # TiB to TB conversion
    
    # SQL Servers: 0 (extract database count)
    sql_servers_match = re.search(r'sql\s+servers?:\s*(\d+)', content, re.IGNORECASE)
    if sql_servers_match:
        result["database_count"] = int(sql_servers_match.group(1))
    
    # Fallback patterns
    if not result["vm_count"]:
        vm_match = re.search(r'(\d+)\s*(?:total\s*)?(?:vms?|virtual\s*machines?|servers?)', content, re.IGNORECASE)
        if vm_match:
            result["vm_count"] = int(vm_match.group(1))
    
    if not result["windows_vms"]:
        windows_match = re.search(r'(\d+)\s*windows', content, re.IGNORECASE)
        if windows_match:
            result["windows_vms"] = int(windows_match.group(1))
    
    if not result["linux_vms"]:
        linux_match = re.search(r'(\d+)\s*linux', content, re.IGNORECASE)
        if linux_match:
            result["linux_vms"] = int(linux_match.group(1))
    
    return result


def extract_executive_summary(prs: Presentation) -> Dict:
    """
    Extract Executive Summary from ATX PowerPoint.
    
    Looks for slides with titles like:
    - "Executive Summary"
    - "Summary"
    
    ATX typically has the executive summary on slide 5 with key metrics
    """
    summary_keywords = ["executive summary"]
    
    # Find all slides with "Executive Summary"
    matching_slides = []
    for idx, slide in enumerate(prs.slides):
        slide_text = extract_text_from_slide(slide).lower()
        if any(keyword.lower() in slide_text for keyword in summary_keywords):
            matching_slides.append((idx, extract_text_from_slide(slide)))
    
    if not matching_slides:
        return {"found": False, "content": ""}
    
    # Use the slide with the most content (usually the detailed one, not just the title slide)
    best_slide = max(matching_slides, key=lambda x: len(x[1]))
    content = best_slide[1]
    
    return {
        "found": True,
        "content": content
    }


def extract_financial_overview(prs: Presentation) -> Dict:
    """
    Extract Financial Overview from ATX PowerPoint.
    
    Looks for slides with titles like:
    - "Financial Overview"
    - "Cost Summary"
    - "Financial Summary"
    - Or in Executive Summary slide with "$" and "annualized"
    """
    financial_keywords = ["financial overview", "cost summary", "financial summary", "pricing summary", "assumptions"]
    
    # Try to find dedicated financial slide
    slide_idx = find_slide_by_title(prs, financial_keywords)
    
    # If not found, check Executive Summary slide for financial data
    if slide_idx is None:
        for idx, slide in enumerate(prs.slides):
            text = extract_text_from_slide(slide)
            if "annualized" in text.lower() and "$" in text:
                slide_idx = idx
                break
    
    if slide_idx is None:
        return {"found": False, "content": ""}
    
    slide = prs.slides[slide_idx]
    content = extract_text_from_slide(slide)
    
    # Parse financial metrics
    result = {
        "found": True,
        "content": content,
        "monthly_cost": None,
        "annual_cost": None,
        "three_year_cost": None
    }
    
    # Extract costs using regex - ATX specific patterns
    # $119,044 annualized spend on AWS
    annual_match = re.search(r'\$\s*([\d,]+(?:\.\d{2})?)\s*annualized', content, re.IGNORECASE)
    if annual_match:
        result["annual_cost"] = float(annual_match.group(1).replace(',', ''))
        result["monthly_cost"] = result["annual_cost"] / 12
        result["three_year_cost"] = result["annual_cost"] * 3
    
    # Fallback patterns
    if not result["annual_cost"]:
        annual_match = re.search(r'\$\s*([\d,]+(?:\.\d{2})?)\s*(?:per\s*)?(?:year|annual)', content, re.IGNORECASE)
        if annual_match:
            result["annual_cost"] = float(annual_match.group(1).replace(',', ''))
            result["monthly_cost"] = result["annual_cost"] / 12
    
    if not result["monthly_cost"]:
        monthly_match = re.search(r'\$\s*([\d,]+(?:\.\d{2})?)\s*(?:per\s*)?month', content, re.IGNORECASE)
        if monthly_match:
            result["monthly_cost"] = float(monthly_match.group(1).replace(',', ''))
    
    if not result["three_year_cost"] and result["annual_cost"]:
        result["three_year_cost"] = result["annual_cost"] * 3
    
    return result


def extract_atx_ppt_data(ppt_file_path: str) -> Dict:
    """
    Main function to extract all key data from ATX PowerPoint.
    
    Args:
        ppt_file_path: Path to ATX PowerPoint file
    
    Returns:
        Dictionary with extracted data from all key slides
    """
    try:
        prs = Presentation(ppt_file_path)
        
        assessment_scope = extract_assessment_scope(prs)
        executive_summary = extract_executive_summary(prs)
        financial_overview = extract_financial_overview(prs)
        
        # If financial data not found in dedicated slide, try Executive Summary
        if financial_overview["found"] and not financial_overview.get("annual_cost"):
            exec_content = executive_summary.get("content", "")
            annual_match = re.search(r'\$\s*([\d,]+(?:\.\d{2})?)\s*annualized', exec_content, re.IGNORECASE)
            if annual_match:
                financial_overview["annual_cost"] = float(annual_match.group(1).replace(',', ''))
                financial_overview["monthly_cost"] = financial_overview["annual_cost"] / 12
                financial_overview["three_year_cost"] = financial_overview["annual_cost"] * 3
                financial_overview["content"] = exec_content  # Use exec summary content for financial
        
        return {
            "success": True,
            "assessment_scope": assessment_scope,
            "executive_summary": executive_summary,
            "financial_overview": financial_overview,
            "total_slides": len(prs.slides)
        }
    
    except Exception as e:
        import traceback
        return {
            "success": False,
            "error": str(e),
            "traceback": traceback.format_exc()
        }


def format_atx_ppt_summary(atx_ppt_data: Dict) -> str:
    """
    Format extracted ATX PowerPoint data into readable summary.
    
    Args:
        atx_ppt_data: Dictionary from extract_atx_ppt_data()
    
    Returns:
        Formatted string summary
    """
    if not atx_ppt_data.get("success"):
        return f"ATX PowerPoint Extraction Failed: {atx_ppt_data.get('error', 'Unknown error')}"
    
    summary = []
    summary.append("=" * 80)
    summary.append("ATX POWERPOINT EXTRACTION SUMMARY")
    summary.append("=" * 80)
    summary.append(f"Total Slides: {atx_ppt_data['total_slides']}")
    summary.append("")
    
    # Assessment Scope
    scope = atx_ppt_data["assessment_scope"]
    if scope["found"]:
        summary.append("ASSESSMENT SCOPE")
        summary.append("-" * 80)
        if scope["vm_count"]:
            summary.append(f"Total VMs: {scope['vm_count']}")
        if scope["vcpu_count"]:
            summary.append(f"Total vCPUs: {scope['vcpu_count']:,}")
        if scope["ram_gb"]:
            summary.append(f"Total RAM: {scope['ram_gb']:,} GB")
        if scope["storage_tb"]:
            summary.append(f"Total Storage: {scope['storage_tb']} TB")
        if scope["windows_vms"]:
            summary.append(f"Windows VMs: {scope['windows_vms']}")
        if scope["linux_vms"]:
            summary.append(f"Linux VMs: {scope['linux_vms']}")
        summary.append("")
    
    # Financial Overview
    financial = atx_ppt_data["financial_overview"]
    if financial["found"]:
        summary.append("FINANCIAL OVERVIEW")
        summary.append("-" * 80)
        if financial["monthly_cost"]:
            summary.append(f"Monthly Cost: ${financial['monthly_cost']:,.2f}")
        if financial["annual_cost"]:
            summary.append(f"Annual Cost: ${financial['annual_cost']:,.2f}")
        if financial["three_year_cost"]:
            summary.append(f"3-Year Cost: ${financial['three_year_cost']:,.2f}")
        summary.append("")
    
    # Executive Summary
    exec_summary = atx_ppt_data["executive_summary"]
    if exec_summary["found"]:
        summary.append("EXECUTIVE SUMMARY")
        summary.append("-" * 80)
        summary.append(exec_summary["content"][:500] + "..." if len(exec_summary["content"]) > 500 else exec_summary["content"])
        summary.append("")
    
    summary.append("=" * 80)
    
    return "\n".join(summary)
